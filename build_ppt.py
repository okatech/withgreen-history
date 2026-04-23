# -*- coding: utf-8 -*-
"""
WithGreen History PPT builder
- 公式コーポレートカラー #006423 を厳守
- ヒーローブラッシュライン（緑×赤・50/10th）モチーフを再現
- 公式ロゴをタイトル/フッターに使用
- フォント: 日本語=Noto Sans JP / 英=DM Sans (システムにあれば、無ければデフォルト)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
import os

from pptx.dml.color import RGBColor

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
OUT    = os.path.join(os.path.dirname(os.path.abspath(__file__)), "WithGreen_History.pptx")

# ---- Brand tokens ----
WG_GREEN      = RGBColor(0x00, 0x64, 0x23)
WG_GREEN_DARK = RGBColor(0x00, 0x4a, 0x1a)
WG_RED        = RGBColor(0xd8, 0x23, 0x2a)
WG_BG         = RGBColor(0xf4, 0xf8, 0xf3)
WG_INK        = RGBColor(0x10, 0x14, 0x10)
WG_MUTED      = RGBColor(0x5a, 0x65, 0x5c)
WHITE         = RGBColor(0xff, 0xff, 0xff)

JP = "Noto Sans JP"
EN = "DM Sans"

# 16:9, EMU
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=WG_INK,
             font=JP, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_brush_lines(slide, cx, cy, total_w):
    """公式signboardの 50 STORES / 10th ANNIVERSARY ブラッシュラインを再現"""
    h = Inches(0.22)
    side_w = (total_w - Inches(2.6)) / 2
    # 緑ライン左
    g = add_rect(slide, cx - total_w/2, cy, side_w, h, WG_GREEN)
    # 50
    add_text(slide, cx - Inches(1.30), cy - Inches(0.45), Inches(1.0), Inches(0.85),
             "50", size=42, bold=True, color=WG_GREEN, font=EN, align=PP_ALIGN.CENTER)
    # 10th
    add_text(slide, cx - Inches(0.10), cy - Inches(0.45), Inches(1.2), Inches(0.85),
             "10th", size=42, bold=True, color=WG_RED, font=EN, align=PP_ALIGN.CENTER)
    # 赤ライン右
    add_rect(slide, cx + Inches(1.3), cy, side_w, h, WG_RED)
    # キャプション
    add_text(slide, cx - Inches(1.30), cy + Inches(0.32), Inches(1.0), Inches(0.30),
             "STORES", size=8, bold=True, color=WG_GREEN, font=EN, align=PP_ALIGN.CENTER)
    add_text(slide, cx - Inches(0.10), cy + Inches(0.32), Inches(1.2), Inches(0.30),
             "ANNIVERSARY", size=8, bold=True, color=WG_RED, font=EN, align=PP_ALIGN.CENTER)


def add_logo(slide, x, y, height=Inches(0.55)):
    p = os.path.join(ASSETS, "wg_logo.png")
    if os.path.exists(p):
        slide.shapes.add_picture(p, x, y, height=height)


def add_footer(slide, page_label):
    bar = add_rect(slide, 0, SH - Inches(0.42), SW, Inches(0.42), WG_GREEN)
    add_text(slide, Inches(0.4), SH - Inches(0.40), Inches(8), Inches(0.4),
             "WITHGREEN HISTORY  ／  創業者・武文智洋社長の言葉でたどる10年の歩み",
             size=9, bold=True, color=WHITE, font=JP, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(2.0), SH - Inches(0.40), Inches(1.6), Inches(0.4),
             page_label, size=9, bold=True, color=WHITE, font=EN,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_top_band(slide, label_en, label_jp):
    """各スライド上部の小さなチャプターラベル"""
    add_rect(slide, 0, 0, SW, Inches(0.35), WG_BG)
    add_text(slide, Inches(0.5), Inches(0.05), Inches(7), Inches(0.28),
             label_en, size=9, bold=True, color=WG_GREEN, font=EN, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(8) - Inches(0.5), Inches(0.05), Inches(8), Inches(0.28),
             label_jp, size=9, bold=True, color=WG_INK, font=JP,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 1 : COVER
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SW, SH, WHITE)

# 左：テキスト塊
add_logo(s, Inches(0.6), Inches(0.5), height=Inches(0.7))

add_text(s, Inches(0.6), Inches(1.6), Inches(7), Inches(0.4),
         "WITHGREEN HISTORY  2016 → 2026",
         size=11, bold=True, color=WG_GREEN, font=EN)

add_brush_lines(s, Inches(3.5), Inches(2.45), Inches(6.0))

add_text(s, Inches(0.6), Inches(3.05), Inches(8), Inches(2.1),
         "サラダボウルで、\n日本の食卓を変える。",
         size=44, bold=True, color=WG_INK, font=JP, spacing=1.25)

add_text(s, Inches(0.6), Inches(5.0), Inches(7.5), Inches(1.2),
         "創業者・武文智洋（たけふみ ともひろ）社長の実際の発言だけで振り返る、\n"
         "サラダボウル専門店『WithGreen』創業から10年・全国50店舗達成までの歩み。",
         size=12, color=WG_MUTED, font=JP, spacing=1.55)

# 右：ポートレート（円形マスク）
portrait = os.path.join(ASSETS, "takefumi-portrait.jpg")
if os.path.exists(portrait):
    pic = s.shapes.add_picture(portrait, Inches(8.9), Inches(1.6), Inches(3.8), Inches(3.8))
    # 円形にトリム
    pic.crop_left = pic.crop_right = 0
    # 円形マスク用の上にOvalを重ねるのは難しいので、別途 oval を背面に置く
# 円形フレーム
frame = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.85), Inches(1.55),
                            Inches(3.9), Inches(3.9))
frame.fill.background()
frame.line.color.rgb = WG_GREEN
frame.line.width = Pt(3)

add_text(s, Inches(8.9), Inches(5.55), Inches(3.8), Inches(0.4),
         "武文 智洋 ／ 代表取締役兼CEO", size=13, bold=True, color=WG_INK,
         font=JP, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.9), Inches(5.95), Inches(3.8), Inches(0.3),
         "TOMOHIRO TAKEFUMI · FOUNDER & CEO", size=9, bold=True,
         color=WG_GREEN, font=EN, align=PP_ALIGN.CENTER)

add_footer(s, "01 / 06")

# ============================================================
# SLIDE 2 : OVERVIEW (店舗数 KPI)
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SW, SH, WHITE)
add_top_band(s, "CHAPTER 01  ／  THE NUMBERS", "数字で見る WithGreen の10年")

# タイトル
add_text(s, Inches(0.5), Inches(0.65), Inches(12), Inches(0.7),
         "創業10年で、店舗数は ", size=28, bold=True, color=WG_INK, font=JP)
# accent
add_text(s, Inches(5.6), Inches(0.65), Inches(2.5), Inches(0.7),
         "1 → 50 店舗", size=28, bold=True, color=WG_GREEN, font=JP)
add_text(s, Inches(8.6), Inches(0.65), Inches(4), Inches(0.7),
         "へ。", size=28, bold=True, color=WG_INK, font=JP)

# KPIカード × 4
kpis = [
    ("2016", "神楽坂で1号店オープン"),
    ("50", "店舗（2026年3月14日達成）"),
    ("2,500,000", "食 ／ 年間サラダ販売数"),
    ("100%", "国産野菜（創業以来）"),
]
card_w = Inches(2.85); card_h = Inches(1.9); gap = Inches(0.2)
total_w = card_w * 4 + gap * 3
start_x = (SW - total_w) / 2
for i, (num, lbl) in enumerate(kpis):
    x = start_x + (card_w + gap) * i
    y = Inches(1.7)
    add_rect(s, x, y, card_w, card_h, WG_BG, line=None)
    add_rect(s, x, y, Inches(0.12), card_h, WG_GREEN)
    add_text(s, x + Inches(0.25), y + Inches(0.25), card_w - Inches(0.4), Inches(0.9),
             num, size=32, bold=True, color=WG_GREEN, font=EN)
    add_text(s, x + Inches(0.25), y + Inches(1.15), card_w - Inches(0.4), Inches(0.7),
             lbl, size=10, bold=True, color=WG_INK, font=JP, spacing=1.4)

# 簡易チャート（バー）
chart_top = Inches(4.0)
add_text(s, Inches(0.5), chart_top - Inches(0.05), Inches(8), Inches(0.4),
         "店舗数の推移（公表値）", size=13, bold=True, color=WG_INK, font=JP)

bars = [("2016", 1), ("2023", 17), ("2024", 30), ("2026", 50)]
chart_y = Inches(4.55); chart_h = Inches(2.2)
chart_x = Inches(1.0); chart_w = Inches(11.3)
# 軸線
add_rect(s, chart_x, chart_y + chart_h, chart_w, Emu(9525), WG_GREEN)
bar_w = Inches(1.6); slot_w = chart_w / len(bars)
max_val = 50
for i, (yr, v) in enumerate(bars):
    bx = chart_x + slot_w * i + (slot_w - bar_w) / 2
    bh = chart_h * (v / max_val)
    by = chart_y + chart_h - bh
    add_rect(s, bx, by, bar_w, bh, WG_GREEN)
    add_text(s, bx - Inches(0.2), by - Inches(0.5), bar_w + Inches(0.4), Inches(0.4),
             f"{v}", size=18, bold=True, color=WG_GREEN, font=EN, align=PP_ALIGN.CENTER)
    add_text(s, bx - Inches(0.2), chart_y + chart_h + Inches(0.05),
             bar_w + Inches(0.4), Inches(0.35),
             yr, size=12, bold=True, color=WG_INK, font=EN, align=PP_ALIGN.CENTER)

add_footer(s, "02 / 06")

# ============================================================
# SLIDE 3 : ORIGIN (2008-2014)
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SW, SH, WHITE)
add_top_band(s, "CHAPTER 02  ／  ORIGIN  2008 - 2014", "原点 ─ ウォール街でサラダボウルに出会う")

add_text(s, Inches(0.5), Inches(0.65), Inches(12), Inches(0.7),
         '「"食"でいこう」とひらめいた瞬間でした。',
         size=26, bold=True, color=WG_INK, font=JP)

# 引用カード3枚（縦並び左カラム）
quotes = [
    ("2008.09  リーマン・ショック",
     "内定先だったリーマン・ブラザーズが経営破綻したことをテレビのニュースで知り、頭が真っ白になり現実を受け止めきれず、その夜は倒れるように寝たのを憶えています。"),
    ("2012  NY・ウォール街",
     "忙しく働く人たちが摂るウォール街でのランチは、肉や穀物がたくさん入った主菜として食べる「サラダボウル」が人気でした。当時の日本には、サラダがメインになる食文化はなく、周りの同僚が毎日食べているのを見てとても驚きました。"),
    ("2014  30歳で退職、共同創業を決意",
     '弟のバックグラウンドにすでに"食"があったことや、これまで自分がやってきたこと、考えてきたこと、未来へのイメージ・・・たくさんのピースがカチリと音を立てたのが聞こえました。「"食"でいこう」とひらめいた瞬間でした。'),
]
qx = Inches(0.5); qy0 = Inches(1.7); qw = Inches(8.4); qh = Inches(1.55); qgap = Inches(0.12)
for i, (t, q) in enumerate(quotes):
    y = qy0 + (qh + qgap) * i
    add_rect(s, qx, y, Inches(0.08), qh, WG_GREEN)
    add_rect(s, qx + Inches(0.08), y, qw - Inches(0.08), qh, WG_BG)
    add_text(s, qx + Inches(0.3), y + Inches(0.12), qw - Inches(0.5), Inches(0.4),
             t, size=12, bold=True, color=WG_GREEN, font=JP)
    add_text(s, qx + Inches(0.3), y + Inches(0.5), qw - Inches(0.5), qh - Inches(0.6),
             "「" + q + "」", size=10, color=WG_INK, font=JP, spacing=1.4)

# 右：商品画像
sal = os.path.join(ASSETS, "salada-mushidori.png")
if os.path.exists(sal):
    s.shapes.add_picture(sal, Inches(9.2), Inches(2.0), height=Inches(3.6))
add_text(s, Inches(9.0), Inches(5.7), Inches(3.8), Inches(0.4),
         "STANDARD SALAD", size=10, bold=True, color=WG_GREEN, font=EN, align=PP_ALIGN.CENTER)
add_text(s, Inches(9.0), Inches(6.0), Inches(3.8), Inches(0.4),
         "蒸し鶏のサラダボウル", size=12, bold=True, color=WG_INK, font=JP, align=PP_ALIGN.CENTER)

add_footer(s, "03 / 06")

# ============================================================
# SLIDE 4 : FOUNDING (2015-2016)
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SW, SH, WHITE)
add_top_band(s, "CHAPTER 03  ／  FOUNDING  2015 - 2016", "創業 ─ 神楽坂で1号店オープン")

add_text(s, Inches(0.5), Inches(0.65), Inches(12), Inches(0.7),
         "「生産者と消費者をつなぐ」サラダボウル専門店、誕生。",
         size=22, bold=True, color=WG_INK, font=JP)

# 大きな引用 + 写真
qbox_x = Inches(0.5); qbox_y = Inches(1.7); qbox_w = Inches(7.5); qbox_h = Inches(4.6)
add_rect(s, qbox_x, qbox_y, Inches(0.1), qbox_h, WG_GREEN)
add_rect(s, qbox_x + Inches(0.1), qbox_y, qbox_w - Inches(0.1), qbox_h, WG_BG)

add_text(s, qbox_x + Inches(0.45), qbox_y + Inches(0.3), Inches(2.5), Inches(0.4),
         "2015 ─ 200日・30カ国の世界一周",
         size=11, bold=True, color=WG_GREEN, font=JP)
add_text(s, qbox_x + Inches(0.45), qbox_y + Inches(0.7), qbox_w - Inches(0.7), Inches(1.7),
         "「日本の良さは、肉・魚・野菜と食材が豊富にあること、料理のアレンジ力もずば抜けて高く、おいしくてコスパがいいこと、四方を海に囲まれ、山や川があり、四季があり、圧倒的な地理の利があります。日本の食文化の豊かさは、世界の中でも稀有です。」",
         size=11, color=WG_INK, font=JP, spacing=1.5)

add_text(s, qbox_x + Inches(0.45), qbox_y + Inches(2.5), Inches(4), Inches(0.4),
         "2016.05 ─ 神楽坂で1号店オープン",
         size=11, bold=True, color=WG_GREEN, font=JP)
add_text(s, qbox_x + Inches(0.45), qbox_y + Inches(2.9), qbox_w - Inches(0.7), Inches(1.6),
         "「アメリカやヨーロッパのサラダ文化をそのまま持ってくるのではなく、『国内の生産者や季節感を大切にするサラダボウル専門店でありたい』という想いで創業しました。だからこそ、創業から大切にしている国産野菜100%の使用を、自然災害などの特殊な事情がない限り、続けていきます。」",
         size=11, color=WG_INK, font=JP, spacing=1.5)

# 右：店舗写真
store = os.path.join(ASSETS, "store_jpg.jpg")
hero  = os.path.join(ASSETS, "hero-salad.jpg")
if os.path.exists(hero):
    s.shapes.add_picture(hero, Inches(8.4), Inches(1.7), width=Inches(4.5), height=Inches(2.8))
if os.path.exists(store):
    s.shapes.add_picture(store, Inches(8.4), Inches(4.6), width=Inches(2.15), height=Inches(1.7))
sm = os.path.join(ASSETS, "salada-roastpork.png")
if os.path.exists(sm):
    s.shapes.add_picture(sm, Inches(10.7), Inches(4.55), width=Inches(2.2), height=Inches(1.8))

add_footer(s, "04 / 06")

# ============================================================
# SLIDE 5 : GROWTH (2023-2026) チャート＋発言
# ============================================================
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SW, SH, WHITE)
add_top_band(s, "CHAPTER 04  ／  GROWTH  2023 - 2026", "成長 ─ 17店舗から50店舗へ")

add_text(s, Inches(0.5), Inches(0.65), Inches(12), Inches(0.7),
         "「日々の選択肢に『サラダボウル』がある未来を創造していきたい。」",
         size=20, bold=True, color=WG_INK, font=JP)

# 左：成長チャート
chart_x = Inches(0.5); chart_y = Inches(2.0); chart_w = Inches(7.5); chart_h = Inches(3.8)
add_rect(s, chart_x, chart_y, chart_w, chart_h, WG_BG)
add_text(s, chart_x + Inches(0.3), chart_y + Inches(0.2), Inches(7), Inches(0.4),
         "店舗数の推移", size=12, bold=True, color=WG_INK, font=JP)
inner_top = chart_y + Inches(0.8); inner_h = chart_h - Inches(1.4)
inner_x = chart_x + Inches(0.5); inner_w = chart_w - Inches(1.0)
add_rect(s, inner_x, inner_top + inner_h, inner_w, Emu(9525), WG_GREEN)
bars2 = [("2016", 1), ("2023", 17), ("2024", 30), ("2026", 50)]
bw = Inches(0.9); slot = inner_w / len(bars2)
for i, (yr, v) in enumerate(bars2):
    bx = inner_x + slot * i + (slot - bw) / 2
    bh = inner_h * (v / 50)
    by = inner_top + inner_h - bh
    add_rect(s, bx, by, bw, bh, WG_GREEN)
    add_text(s, bx - Inches(0.3), by - Inches(0.45), bw + Inches(0.6), Inches(0.4),
             str(v), size=14, bold=True, color=WG_GREEN, font=EN, align=PP_ALIGN.CENTER)
    add_text(s, bx - Inches(0.3), inner_top + inner_h + Inches(0.05),
             bw + Inches(0.6), Inches(0.3),
             yr, size=10, bold=True, color=WG_INK, font=EN, align=PP_ALIGN.CENTER)

# 右：発言カード
qx = Inches(8.4); qy = Inches(2.0); qw = Inches(4.5)
quotes2 = [
    ("2023  ビジョン",
     "５年後には30〜50店舗に拡大し、日々の選択肢に『サラダボウル』がある未来を創造していきたい。"),
    ("2026.03.14  創業10年・50店舗達成",
     "ここから始まる2026年からの10年間では、さらなる店舗拡大を目指し、『日本を代表するサラダボウル専門店』へと邁進いたします。"),
]
for i, (t, q) in enumerate(quotes2):
    y = qy + Inches(2.0) * i
    add_rect(s, qx, y, Inches(0.08), Inches(1.85), WG_GREEN)
    add_rect(s, qx + Inches(0.08), y, qw - Inches(0.08), Inches(1.85), WG_BG)
    add_text(s, qx + Inches(0.25), y + Inches(0.12), qw - Inches(0.4), Inches(0.4),
             t, size=11, bold=True, color=WG_GREEN, font=JP)
    add_text(s, qx + Inches(0.25), y + Inches(0.5), qw - Inches(0.4), Inches(1.3),
             "「" + q + "」", size=10, color=WG_INK, font=JP, spacing=1.5)

add_footer(s, "05 / 06")

# ============================================================
# SLIDE 6 : CLOSING
# ============================================================
s = prs.slides.add_slide(blank)
# 全面 WG_GREEN
add_rect(s, 0, 0, SW, SH, WG_GREEN)

add_logo(s, Inches(0.6), Inches(0.5), height=Inches(0.7))
# ロゴを白く反転できないので別表現：白テキストのウォードマーク
add_text(s, Inches(1.6), Inches(0.55), Inches(8), Inches(0.6),
         "WITHGREEN", size=22, bold=True, color=WHITE, font=EN)

# 中央メッセージ
add_brush_lines(s, SW/2, Inches(2.4), Inches(7.2))

add_text(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.6),
         "サラダボウルを、\n日本の新しい食文化へ。",
         size=46, bold=True, color=WHITE, font=JP, align=PP_ALIGN.CENTER, spacing=1.25)

add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.6),
         "WITHGREEN  WITH  YOUR  LIFE",
         size=14, bold=True, color=WHITE, font=EN, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
         "出典：株式会社ネタもと『広報PRのチカラ』(2023) ／ PR TIMESプレスリリース (2026.3) ／ withgreen.club",
         size=9, color=WHITE, font=JP, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
         "本資料はファンメイドの非公式アーカイブです。LOGO © WITHGREEN, INC.",
         size=9, color=WHITE, font=JP, align=PP_ALIGN.CENTER)

# フッター（白いライン）
add_rect(s, 0, SH - Inches(0.05), SW, Inches(0.05), WHITE)

prs.save(OUT)
print("OK ->", OUT)
